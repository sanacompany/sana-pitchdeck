#!/usr/bin/env python3
"""Extract PPTX content, handling invalid XML characters."""
import json, os, sys, re, zipfile, shutil, tempfile

def sanitize_xml_in_pptx(input_path):
    """Re-pack the PPTX with invalid XML chars stripped from XML parts."""
    tmp = tempfile.mkdtemp()
    clean_path = os.path.join(tmp, "clean.pptx")
    invalid_xml_re = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f]')

    with zipfile.ZipFile(input_path, 'r') as zin, zipfile.ZipFile(clean_path, 'w') as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(('.xml', '.rels')):
                text = data.decode('utf-8', errors='replace')
                text = invalid_xml_re.sub('', text)
                data = text.encode('utf-8')
            zout.writestr(item, data)
    return clean_path, tmp

def extract_pptx(file_path, output_dir="."):
    from pptx import Presentation

    clean_path, tmp_dir = sanitize_xml_in_pptx(file_path)
    try:
        prs = Presentation(clean_path)
        slides_data = []
        assets_dir = os.path.join(output_dir, "assets")
        os.makedirs(assets_dir, exist_ok=True)

        for slide_num, slide in enumerate(prs.slides):
            slide_data = {
                "number": slide_num + 1,
                "title": "",
                "content": [],
                "images": [],
                "notes": "",
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape == slide.shapes.title:
                        slide_data["title"] = shape.text
                    else:
                        slide_data["content"].append({"type": "text", "content": shape.text})
                if shape.shape_type == 13:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                    image_path = os.path.join(assets_dir, image_name)
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    slide_data["images"].append({
                        "path": f"assets/{image_name}",
                        "width": shape.width,
                        "height": shape.height,
                    })
            if slide.has_notes_slide:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text
            slides_data.append(slide_data)
        return slides_data
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

if __name__ == "__main__":
    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."
    os.makedirs(output_dir, exist_ok=True)
    slides = extract_pptx(input_file, output_dir)
    output_path = os.path.join(output_dir, "extracted-slides.json")
    with open(output_path, "w") as f:
        json.dump(slides, f, indent=2)
    print(f"Extracted {len(slides)} slides to {output_path}")
    for s in slides:
        img_count = len(s["images"])
        print(f"  Slide {s['number']}: {s['title'] or '(no title)'} — {img_count} image(s)")
